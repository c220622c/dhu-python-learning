#!/usr/bin/env python3
from pptx import Presentation
import sys

def read_pptx_outline(file_path):
    """Extract outline from PowerPoint file"""
    try:
        prs = Presentation(file_path)
        outline = []

        for i, slide in enumerate(prs.slides):
            slide_content = f"Slide {i + 1}:"

            # Get slide title
            if slide.shapes.title:
                slide_content += f"\n  Title: {slide.shapes.title.text}"

            # Get text from other shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
                    # Only include non-empty text that's not the title
                    text = shape.text.strip()
                    if text and len(text) > 10:  # Skip very short text fragments
                        slide_content += f"\n  • {text[:100]}{'...' if len(text) > 100 else ''}"

            outline.append(slide_content)

        return outline

    except Exception as e:
        return [f"Error reading {file_path}: {str(e)}"]

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python3 read_pptx.py <pptx_file>")
        sys.exit(1)

    file_path = sys.argv[1]
    outline = read_pptx_outline(file_path)

    for slide in outline:
        print(slide)
        print()